# -*- coding: utf-8 -*-
"""
Google Drive API helper — OAuth2 authentication + file search + content reading.

First-time use:
  1. Enable the Google Drive API in Google Cloud Console (same project as Gmail/Calendar).
  2. The existing credentials.json already works — run the agent and ask about Drive.
     A browser window opens for OAuth consent.
  3. The token is saved to gdrive_token.json for future runs.

Scopes: drive.readonly (search + read only — no writes)

Supports reading:
  - Google Docs        → plain text
  - Google Sheets      → CSV text
  - Google Slides      → plain text
  - PDF files          → plain text (via export)
  - Plain text / code  → raw content
  - Other formats      → metadata only
"""
from __future__ import annotations

import io
import os
from dataclasses import dataclass, field
from typing import List, Optional

_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

_DEFAULT_CREDENTIALS = os.path.join(os.path.dirname(__file__), "credentials.json")
_DEFAULT_TOKEN = os.path.join(os.path.dirname(__file__), "gdrive_token.json")

# MIME types we can extract text from
_EXPORTABLE = {
    "application/vnd.google-apps.document":     ("text/plain",    "Google Doc"),
    "application/vnd.google-apps.spreadsheet":  ("text/csv",      "Google Sheet"),
    "application/vnd.google-apps.presentation": ("text/plain",    "Google Slides"),
}
_PLAIN_TEXT_MIMES = {
    "text/plain", "text/markdown", "text/csv", "text/html",
    "application/json", "text/javascript", "application/x-python",
}
# Office / binary formats we can parse locally
_DOCX_MIME  = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_DOC_MIME   = "application/msword"
_XLSX_MIME  = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_XLS_MIME   = "application/vnd.ms-excel"
_PPTX_MIME  = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_OFFICE_MIMES = {_DOCX_MIME, _DOC_MIME, _XLSX_MIME, _XLS_MIME, _PPTX_MIME}
_PDF_MIME = "application/pdf"

MAX_CONTENT_CHARS = 6000   # cap returned text to avoid context overflow


# ---------------------------------------------------------------------------
# Auth
# ---------------------------------------------------------------------------

def _get_service(
    credentials_path: str = _DEFAULT_CREDENTIALS,
    token_path: str = _DEFAULT_TOKEN,
):
    from google.auth.transport.requests import Request
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from googleapiclient.discovery import build

    creds = None
    if os.path.exists(token_path):
        creds = Credentials.from_authorized_user_file(token_path, _SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not os.path.exists(credentials_path):
                raise FileNotFoundError(
                    f"credentials.json not found at {credentials_path}."
                )
            flow = InstalledAppFlow.from_client_secrets_file(credentials_path, _SCOPES)
            creds = flow.run_local_server(port=0)
        with open(token_path, "w") as fh:
            fh.write(creds.to_json())

    return build("drive", "v3", credentials=creds)


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class DriveFile:
    file_id: str
    name: str
    mime_type: str
    modified: str = ""
    size: str = ""
    owners: List[str] = field(default_factory=list)
    web_link: str = ""
    kind: str = ""          # human label: "Google Doc", "PDF", etc.


# ---------------------------------------------------------------------------
# Search
# ---------------------------------------------------------------------------

def search_files(
    query: str,
    max_results: int = 10,
    credentials_path: str = _DEFAULT_CREDENTIALS,
    token_path: str = _DEFAULT_TOKEN,
) -> List[DriveFile]:
    """
    Search Google Drive for files matching `query`.
    `query` is free-text: file names, keywords, or Drive query syntax.
    Returns a list of DriveFile objects.
    """
    service = _get_service(credentials_path, token_path)

    # Build Drive API query: full-text search OR name match
    # Drive query syntax: https://developers.google.com/drive/api/guides/search-files
    if any(op in query for op in ("=", "contains", "mimeType", "in parents", "starred")):
        drive_q = query  # caller passed raw Drive query syntax
    else:
        # Free-text: search name and full text
        safe = query.replace("'", "\\'")
        drive_q = f"fullText contains '{safe}' and trashed = false"

    fields = "files(id,name,mimeType,modifiedTime,size,owners,webViewLink)"
    result = (
        service.files()
        .list(
            q=drive_q,
            pageSize=max_results,
            fields=fields,
            orderBy="modifiedTime desc",
        )
        .execute()
    )

    files = []
    for item in result.get("files", []):
        mime = item.get("mimeType", "")
        kind = _EXPORTABLE.get(mime, ("", mime.split("/")[-1]))[1]
        if mime == _PDF_MIME:
            kind = "PDF"
        files.append(DriveFile(
            file_id=item.get("id", ""),
            name=item.get("name", "(unnamed)"),
            mime_type=mime,
            modified=item.get("modifiedTime", "")[:10],
            size=str(item.get("size", "")),
            owners=[o.get("displayName", "") for o in item.get("owners", [])],
            web_link=item.get("webViewLink", ""),
            kind=kind,
        ))
    return files


def format_search_results(files: List[DriveFile], query: str) -> str:
    if not files:
        return f"🔍 No Drive files found for: {query!r}"

    lines = [f"🔍 **Drive search: '{query}'** — {len(files)} result(s)\n"]
    for i, f in enumerate(files, 1):
        size_str = f"  {int(f.size)//1024} KB" if f.size.isdigit() else ""
        lines.append(
            f"  {i}. **{f.name}** [{f.kind}]{size_str}\n"
            f"     📅 Modified: {f.modified}  |  🆔 `{f.file_id}`\n"
            f"     🔗 {f.web_link}"
        )
    lines.append("\n💡 To read a file: 读取文件 <file_id> or 'read file 1' (by result number)")
    return "\n".join(lines)


def search_files_as_text(
    query: str,
    max_results: int = 10,
    credentials_path: str = _DEFAULT_CREDENTIALS,
    token_path: str = _DEFAULT_TOKEN,
) -> str:
    try:
        files = search_files(query, max_results, credentials_path, token_path)
        return format_search_results(files, query)
    except FileNotFoundError as e:
        return f"❌ Drive setup error: {e}"
    except Exception as e:
        return f"❌ Drive search failed: {e}"


# ---------------------------------------------------------------------------
# Read file content
# ---------------------------------------------------------------------------

def read_file_content(
    file_id: str,
    max_chars: int = MAX_CONTENT_CHARS,
    credentials_path: str = _DEFAULT_CREDENTIALS,
    token_path: str = _DEFAULT_TOKEN,
) -> str:
    """
    Read the text content of a Drive file by its file_id.
    Returns formatted text with metadata header + content.
    """
    service = _get_service(credentials_path, token_path)

    # Get file metadata
    meta = service.files().get(
        fileId=file_id,
        fields="id,name,mimeType,modifiedTime,owners,webViewLink"
    ).execute()

    name = meta.get("name", file_id)
    mime = meta.get("mimeType", "")
    modified = meta.get("modifiedTime", "")[:10]
    web_link = meta.get("webViewLink", "")

    header = (
        f"📄 **{name}**\n"
        f"   Type: {mime}  |  Modified: {modified}\n"
        f"   🔗 {web_link}\n"
        f"{'─'*50}\n"
    )

    content = _extract_text(service, file_id, mime, max_chars)
    return header + content


def _extract_text(service, file_id: str, mime: str, max_chars: int) -> str:
    """Extract text content depending on file type."""
    # Google Workspace files (Docs, Sheets, Slides) → export as text
    if mime in _EXPORTABLE:
        export_mime, kind = _EXPORTABLE[mime]
        try:
            data = service.files().export(
                fileId=file_id, mimeType=export_mime
            ).execute()
            if isinstance(data, bytes):
                text = data.decode("utf-8", errors="replace")
            else:
                text = str(data)
            text = text.strip()
            if len(text) > max_chars:
                text = text[:max_chars] + f"\n\n… [truncated — {len(text)} chars total]"
            return text or "(empty document)"
        except Exception as e:
            return f"❌ Could not export {kind}: {e}"

    # PDF → download and extract text with pypdf (or fall back gracefully)
    if mime == _PDF_MIME:
        try:
            data = _download_bytes(service, file_id)
        except Exception as e:
            return f"❌ Could not download PDF: {e}"
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(data))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n\n".join(p.strip() for p in pages if p.strip())
            if len(text) > max_chars:
                text = text[:max_chars] + f"\n\n… [truncated — {len(text)} chars total]"
            return text or "(no extractable text in PDF)"
        except ImportError:
            return (
                "⚠️  pypdf is not installed. Run:\n"
                "    pip install pypdf\n"
                "Then retry."
            )
        except Exception as e:
            return f"❌ Could not parse PDF: {e}"

    # Plain text / code / markdown → download directly
    if mime in _PLAIN_TEXT_MIMES or mime.startswith("text/"):
        try:
            data = _download_bytes(service, file_id)
            text = data.decode("utf-8", errors="replace").strip()
            if len(text) > max_chars:
                text = text[:max_chars] + f"\n\n… [truncated]"
            return text or "(empty file)"
        except Exception as e:
            return f"❌ Could not read file: {e}"

    # Office formats (.docx, .xlsx, .pptx, .doc) → download + parse locally
    if mime in _OFFICE_MIMES:
        try:
            data = _download_bytes(service, file_id)
        except Exception as e:
            return f"❌ Could not download file: {e}"
        if mime in (_DOCX_MIME, _DOC_MIME):
            return _read_docx(data, max_chars)
        if mime in (_XLSX_MIME, _XLS_MIME):
            return _read_xlsx(data, max_chars)
        if mime == _PPTX_MIME:
            return _read_pptx(data, max_chars)

    # Unsupported binary type
    return (
        f"⚠️  Cannot read binary file (MIME: {mime}).\n"
        f"Open it directly in your browser via the link above."
    )


def _download_bytes(service, file_id: str) -> bytes:
    """Download a file's raw binary content."""
    import googleapiclient.http as ghttp
    request = service.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    downloader = ghttp.MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue()


def _read_docx(data: bytes, max_chars: int) -> str:
    try:
        import docx  # python-docx
        doc = docx.Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n… [truncated — {len(text)} chars total]"
        return text or "(empty document)"
    except ImportError:
        return (
            "⚠️  python-docx is not installed. Run:\n"
            "    pip install python-docx\n"
            "Then retry."
        )
    except Exception as e:
        return f"❌ Could not parse .docx: {e}"


def _read_xlsx(data: bytes, max_chars: int) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                row_str = "\t".join("" if v is None else str(v) for v in row)
                if row_str.strip():
                    lines.append(row_str)
        text = "\n".join(lines)
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n… [truncated — {len(text)} chars total]"
        return text or "(empty spreadsheet)"
    except ImportError:
        return (
            "⚠️  openpyxl is not installed. Run:\n"
            "    pip install openpyxl\n"
            "Then retry."
        )
    except Exception as e:
        return f"❌ Could not parse .xlsx: {e}"


def _read_pptx(data: bytes, max_chars: int) -> str:
    try:
        from pptx import Presentation  # python-pptx
        prs = Presentation(io.BytesIO(data))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        text = "\n".join(lines)
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n… [truncated — {len(text)} chars total]"
        return text or "(no text in presentation)"
    except ImportError:
        return (
            "⚠️  python-pptx is not installed. Run:\n"
            "    pip install python-pptx\n"
            "Then retry."
        )
    except Exception as e:
        return f"❌ Could not parse .pptx: {e}"
